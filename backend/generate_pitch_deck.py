"""
Generate Professional Pitch Deck for Ross Lending Solutions LLC
For bank presentations and investor meetings.
Generates both ES and EN versions as .pptx files.

Run: python3 generate_pitch_deck.py
"""
import os
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime


# ═══════════════════════════════════════════════════════════════════════════
# BRAND COLORS
# ═══════════════════════════════════════════════════════════════════════════
EMERALD = RGBColor(0x05, 0x96, 0x69)
DARK_EMERALD = RGBColor(0x04, 0x78, 0x57)
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
DARK_CARD = RGBColor(0x1E, 0x29, 0x3B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
MEDIUM_GRAY = RGBColor(0x64, 0x74, 0x8B)
TEXT_DARK = RGBColor(0x1E, 0x29, 0x3B)
GOLD = RGBColor(0xF5, 0x9E, 0x0B)
RED_ACCENT = RGBColor(0xEF, 0x44, 0x44)
BLUE_ACCENT = RGBColor(0x38, 0xBD, 0xF8)


def add_dark_bg(slide, prs):
    """Add dark background to slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_shape_bg(slide, left, top, width, height, color, corner_radius=None):
    """Add a colored rectangle shape as background."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if corner_radius:
        shape.adjustments[0] = corner_radius
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box with specified properties."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_footer(slide, text="Ross Lending Solutions LLC | CONFIDENCIAL"):
    """Add footer bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.0), Inches(13.33), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_CARD
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.33), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER


def add_table(slide, left, top, width, rows_data, col_widths, header_color=EMERALD):
    """Add a styled table."""
    rows = len(rows_data)
    cols = len(rows_data[0]) if rows_data else 0
    if rows == 0 or cols == 0:
        return

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(rows * 0.45))
    table = table_shape.table

    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)

    for r_idx, row in enumerate(rows_data):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = 'Calibri'
                if r_idx == 0:
                    paragraph.font.color.rgb = WHITE
                    paragraph.font.bold = True
                else:
                    paragraph.font.color.rgb = TEXT_DARK if r_idx % 2 == 1 else RGBColor(0x33, 0x41, 0x55)

            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xFD, 0xF4)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table_shape


def build_presentation(lang='es'):
    """Build the full pitch deck."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    t = TEXTS[lang]

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 1: COVER
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_dark_bg(slide, prs)

    # Emerald accent bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = EMERALD
    bar.line.fill.background()

    # Company name
    add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1),
                 "ROSS LENDING SOLUTIONS LLC", 48, WHITE, True, PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, Inches(1), Inches(2.9), Inches(11), Inches(0.8),
                 t['cover_subtitle'], 24, EMERALD, False, PP_ALIGN.CENTER)

    # Tagline
    add_text_box(slide, Inches(2), Inches(3.8), Inches(9), Inches(0.6),
                 t['cover_tagline'], 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.6), Inches(2.33), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = EMERALD
    line.line.fill.background()

    # Contact info
    add_text_box(slide, Inches(2), Inches(5.0), Inches(9), Inches(1.2),
                 "305 Bruce Ave, Dumas, TX 79029 | (806) 934-2018\ninfo@rosslending.com | www.rosslending.com\n\nTexas OCCC Regulated Lender | Chapter 342 Texas Finance Code",
                 13, MEDIUM_GRAY, False, PP_ALIGN.CENTER)

    # Date
    date_str = datetime.now().strftime('%B %Y') if lang == 'en' else datetime.now().strftime('%B %Y')
    add_text_box(slide, Inches(2), Inches(6.4), Inches(9), Inches(0.4),
                 f"{t['confidential']} | {date_str}", 11, RGBColor(0x47, 0x55, 0x69), False, PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 2: EXECUTIVE SUMMARY
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide, prs)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 t['exec_title'], 32, WHITE, True)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = EMERALD; line.line.fill.background()

    # Three cards
    cards = t['exec_cards']
    for i, card in enumerate(cards):
        x = Inches(0.8 + i * 4)
        bg = add_shape_bg(slide, x, Inches(1.5), Inches(3.6), Inches(4.8), DARK_CARD)

        add_text_box(slide, x + Inches(0.3), Inches(1.7), Inches(3), Inches(0.5),
                     card['icon'], 36, EMERALD, False, PP_ALIGN.LEFT)
        add_text_box(slide, x + Inches(0.3), Inches(2.3), Inches(3), Inches(0.5),
                     card['title'], 18, WHITE, True)

        # Bullet points
        txBox = slide.shapes.add_textbox(x + Inches(0.3), Inches(2.9), Inches(3), Inches(3.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        for j, bullet in enumerate(card['bullets']):
            p = tf.add_paragraph() if j > 0 else tf.paragraphs[0]
            p.text = f"• {bullet}"
            p.font.size = Pt(12)
            p.font.color.rgb = LIGHT_GRAY
            p.font.name = 'Calibri'
            p.space_after = Pt(6)

    add_footer(slide, t['footer'])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 3: REGULATORY FRAMEWORK
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide, prs)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 t['reg_title'], 32, WHITE, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = EMERALD; line.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
                 t['reg_subtitle'], 14, LIGHT_GRAY)

    add_table(slide, Inches(0.8), Inches(2.1), Inches(11.5),
              t['reg_table'], [2.5, 2, 2, 1.8, 1.5, 1.7])

    # Compliance bullets
    add_text_box(slide, Inches(0.8), Inches(4.5), Inches(5), Inches(0.4),
                 t['compliance_title'], 18, EMERALD, True)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(t['compliance_items']):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"✅  {item}"
        p.font.size = Pt(13)
        p.font.color.rgb = LIGHT_GRAY
        p.font.name = 'Calibri'
        p.space_after = Pt(4)

    add_footer(slide, t['footer'])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 4: LOAN PRODUCTS
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide, prs)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 t['products_title'], 32, WHITE, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = EMERALD; line.line.fill.background()

    # Sub F Card
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
                 t['subf_title'], 20, GOLD, True)
    add_table(slide, Inches(0.8), Inches(2.1), Inches(5.5),
              t['subf_table'], [1.0, 0.8, 0.9, 1.0, 0.9], GOLD)

    # Sub E Card
    add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
                 t['sube_title'], 20, EMERALD, True)
    add_table(slide, Inches(7), Inches(2.1), Inches(5.5),
              t['sube_table'], [1.0, 0.8, 0.9, 1.0, 0.9], EMERALD)

    # Tax Advance
    add_shape_bg(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.3), DARK_CARD)
    add_text_box(slide, Inches(1.2), Inches(5.35), Inches(3), Inches(0.4),
                 t['tax_advance_title'], 18, BLUE_ACCENT, True)
    add_text_box(slide, Inches(1.2), Inches(5.8), Inches(10), Inches(0.6),
                 t['tax_advance_desc'], 13, LIGHT_GRAY)

    add_footer(slide, t['footer'])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 5: CAPITAL STRATEGY
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide, prs)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 t['capital_title'], 32, WHITE, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = EMERALD; line.line.fill.background()

    # $27K scenario
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
                 t['scenario_27k'], 18, EMERALD, True)
    add_table(slide, Inches(0.8), Inches(2.1), Inches(5.5),
              t['table_27k'], [1.5, 1.0, 1.2, 1.0], EMERALD)

    # $80K scenario
    add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
                 t['scenario_80k'], 18, GOLD, True)
    add_table(slide, Inches(7), Inches(2.1), Inches(5.5),
              t['table_80k'], [1.5, 1.0, 1.2, 1.0], GOLD)

    # Key metrics
    metrics = t['key_metrics']
    for i, m in enumerate(metrics):
        x = Inches(0.8 + i * 3.1)
        bg = add_shape_bg(slide, x, Inches(4.8), Inches(2.8), Inches(1.8), DARK_CARD)
        add_text_box(slide, x + Inches(0.3), Inches(4.95), Inches(2.2), Inches(0.8),
                     m['value'], 28, EMERALD, True, PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.3), Inches(5.7), Inches(2.2), Inches(0.5),
                     m['label'], 12, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    add_footer(slide, t['footer'])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 6: INVESTOR MODEL
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide, prs)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 t['investor_title'], 32, WHITE, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = EMERALD; line.line.fill.background()

    add_table(slide, Inches(0.8), Inches(1.5), Inches(11.5),
              t['investor_table'], [3, 2, 2, 2, 2.5], EMERALD)

    add_text_box(slide, Inches(0.8), Inches(4.0), Inches(5), Inches(0.4),
                 t['sec_title'], 18, EMERALD, True)
    add_table(slide, Inches(0.8), Inches(4.5), Inches(11.5),
              t['sec_table'], [3, 4, 2, 2.5], DARK_CARD)

    add_footer(slide, t['footer'])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 7: FINANCIAL PROJECTIONS
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide, prs)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 t['projection_title'], 32, WHITE, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = EMERALD; line.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.4),
                 t['projection_subtitle'], 13, LIGHT_GRAY)

    add_table(slide, Inches(0.8), Inches(2.0), Inches(11.5),
              t['projection_table'], [1.5, 2, 2.5, 2, 2.5, 1], EMERALD)

    # Risk management section
    add_text_box(slide, Inches(0.8), Inches(4.8), Inches(5), Inches(0.4),
                 t['risk_title'], 18, RED_ACCENT, True)
    add_table(slide, Inches(0.8), Inches(5.3), Inches(11.5),
              t['risk_table'], [3, 1.5, 1.5, 5.5], RGBColor(0xDC, 0x26, 0x26))

    add_footer(slide, t['footer'])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 8: 90-DAY ACTION PLAN
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide, prs)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 t['action_title'], 32, WHITE, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = EMERALD; line.line.fill.background()

    phases = t['phases']
    colors = [EMERALD, GOLD, BLUE_ACCENT]
    for i, phase in enumerate(phases):
        x = Inches(0.8 + i * 4)
        bg = add_shape_bg(slide, x, Inches(1.5), Inches(3.6), Inches(5.0), DARK_CARD)

        add_text_box(slide, x + Inches(0.3), Inches(1.7), Inches(3), Inches(0.4),
                     phase['label'], 14, colors[i], True)
        add_text_box(slide, x + Inches(0.3), Inches(2.1), Inches(3), Inches(0.5),
                     phase['title'], 20, WHITE, True)

        txBox = slide.shapes.add_textbox(x + Inches(0.3), Inches(2.8), Inches(3), Inches(3.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for j, item in enumerate(phase['items']):
            p = tf.add_paragraph() if j > 0 else tf.paragraphs[0]
            p.text = f"→ {item}"
            p.font.size = Pt(12)
            p.font.color.rgb = LIGHT_GRAY
            p.font.name = 'Calibri'
            p.space_after = Pt(8)

    add_footer(slide, t['footer'])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 9: CLOSING / THANK YOU
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide, prs)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = EMERALD; bar.line.fill.background()

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
                 t['closing_title'], 42, WHITE, True, PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.6),
                 "ROSS LENDING SOLUTIONS LLC", 28, EMERALD, True, PP_ALIGN.CENTER)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.6), Inches(2.33), Inches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = EMERALD; line.line.fill.background()

    contact_lines = [
        "Yoandy Ross — CEO & Founder",
        "(806) 934-2018 | info@rosslending.com",
        "305 Bruce Ave, Dumas, TX 79029",
        "www.rosslending.com",
        "",
        t['closing_license'],
    ]
    txBox = slide.shapes.add_textbox(Inches(2), Inches(3.9), Inches(9), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_text in enumerate(contact_lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line_text
        p.font.size = Pt(15)
        p.font.color.rgb = LIGHT_GRAY
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(4)

    return prs


# ═══════════════════════════════════════════════════════════════════════════
# TEXTS - ES & EN
# ═══════════════════════════════════════════════════════════════════════════

TEXTS = {
    'es': {
        'cover_subtitle': 'Plan de Negocio Integral 2026-2028',
        'cover_tagline': 'Préstamos Inteligentes. Resultados Reales. | Tu Socio Financiero en Cada Paso.',
        'confidential': 'CONFIDENCIAL — Solo para uso interno',
        'footer': 'Ross Lending Solutions LLC | 305 Bruce Ave, Dumas, TX 79029 | CONFIDENCIAL',

        'exec_title': 'Resumen Ejecutivo',
        'exec_cards': [
            {
                'icon': '🏢',
                'title': 'La Empresa',
                'bullets': [
                    'Prestamista regulado en Texas bajo licencia OCCC Cap. 342',
                    'Préstamos personales corto y mediano plazo',
                    'Enfoque: comunidad hispana del Texas Panhandle',
                    'Servicio 100% bilingüe (ES/EN)',
                    'Aprobación en 24 horas',
                ]
            },
            {
                'icon': '🎯',
                'title': 'Oportunidad de Mercado',
                'bullets': [
                    '65% de la población del Panhandle es hispana',
                    'Acceso limitado a crédito bancario tradicional',
                    'Demanda de $2-5M anuales en microcréditos',
                    'Poca competencia regulada en la zona',
                    'Base de 900+ clientes existentes (Ross Tax)',
                ]
            },
            {
                'icon': '💰',
                'title': 'Propuesta de Inversión',
                'bullets': [
                    'Capital inicial: $27,000 - $80,000',
                    'ROI proyectado: 180-240% anual',
                    'Retorno al inversionista: 8-15% anual',
                    'Préstamos diversificados (50+ micro-préstamos)',
                    'Regulado y conforme a ley estatal',
                ]
            },
        ],

        'reg_title': 'Marco Regulatorio',
        'reg_subtitle': 'Licencia OCCC — Texas Finance Code, Chapter 342 (Una sola licencia cubre ambos subcapítulos)',
        'reg_table': [
            ['Subcapítulo', 'Sección Legal', 'Tipo', 'Montos', 'Plazos', 'Tasa Máxima'],
            ['Subcapítulo F', '§342.251-259', 'Corto plazo', '$100 - $1,500', '1-3 meses', '180-240% APR'],
            ['Subcapítulo E', '§342.201-214', 'Plazos / Installment', '$2,000 - $12,000', '6-48 meses', '18-30% APR'],
        ],
        'compliance_title': 'Cumplimiento',
        'compliance_items': [
            'Truth in Lending Act (TILA) — Divulgación obligatoria de APR y cargos',
            'Regulation Z — Formato estándar de divulgación',
            'Equal Credit Opportunity Act — No discriminación',
            'FDCPA & Texas Debt Collection Act — Prácticas justas de cobro',
        ],

        'products_title': 'Productos de Préstamo',
        'subf_title': '⚡ Préstamo Rápido — Subcapítulo F',
        'subf_table': [
            ['Monto', 'APR', 'Plazo', 'Cargo', 'Total'],
            ['$200', '240%', '1 mes', '$40', '$240'],
            ['$500', '180%', '1 mes', '$75', '$575'],
            ['$1,000', '180%', '2 meses', '$170', '$1,170'],
            ['$1,500', '180%', '3 meses', '$240', '$1,740'],
        ],
        'sube_title': '📊 Préstamo a Plazos — Subcapítulo E',
        'sube_table': [
            ['Monto', 'Tasa', 'Plazo', 'Pago/Mes', 'Total'],
            ['$2,000', '~28%', '12 meses', '$130', '$2,425'],
            ['$5,000', '~21%', '24 meses', '$255', '$6,245'],
            ['$8,000', '~20%', '36 meses', '$297', '$10,825'],
            ['$12,000', '~18%', '48 meses', '$353', '$16,940'],
        ],
        'tax_advance_title': '🏛️ Adelanto de Taxes',
        'tax_advance_desc': 'Préstamo puente (Ene-Abr) que se paga automáticamente del refund del IRS. Montos: $200-$1,800 | Sub F rates | Base de 900+ clientes de temporada fiscal.',

        'capital_title': 'Estrategia de Capital',
        'scenario_27k': '📋 Escenario: $27,000 (Híbrido)',
        'table_27k': [
            ['Componente', 'Capital', 'Prést/Mes', 'Ganancia'],
            ['Sub F — 70%', '$18,900', '~38', '$2,850/mes'],
            ['Sub E — 30%', '$8,100', '~3', '$160/mes'],
            ['TOTAL', '$27,000', '~41', '$3,010+/mes'],
        ],
        'scenario_80k': '🚀 Escenario: $80,000 (Escalamiento)',
        'table_80k': [
            ['Componente', 'Capital', 'Prést/Mes', 'Ganancia'],
            ['Sub F Mix', '$48,000', '~96', '$7,688/mes'],
            ['Sub E Mix', '$24,000', '~6', '$362/mes'],
            ['TOTAL', '$80,000', '~100+', '$8,050+/mes'],
        ],
        'key_metrics': [
            {'value': '$3,010+', 'label': 'Ingreso Mensual\n(con $27K)'},
            {'value': '$8,050+', 'label': 'Ingreso Mensual\n(con $80K)'},
            {'value': '180-240%', 'label': 'ROI Anual\nProyectado'},
            {'value': '8-15%', 'label': 'Retorno Anual\nal Inversionista'},
        ],

        'investor_title': 'Modelo de Inversionistas',
        'investor_table': [
            ['Tipo de Nota', 'Retorno Anual', 'Riesgo', 'Inversión Mín.', 'Pago de Intereses'],
            ['Pagaré Privado', '8-10%', 'Bajo', '$10,000', 'Mensual o Trimestral'],
            ['Nota Participativa', '10-12%', 'Moderado', '$25,000', 'Trimestral'],
            ['Participación en Ganancias', '12-15%', 'Medio', '$50,000', 'Mensual'],
            ['Partnership', '15-20%', 'Alto', '$100,000', 'Mensual'],
        ],
        'sec_title': 'Exenciones SEC / Texas Securities',
        'sec_table': [
            ['Exención', 'Requisitos', 'Costo Legal', 'Fase'],
            ['Texas §4005.012(a)(2)', 'Máx 15 inversionistas, sofisticados, solo TX', '$0', 'Fase 1'],
            ['Federal Rule 506(b)', 'Máx 35 no-acreditados, PPM requerido', '$2K-$5K', 'Fase 2'],
            ['Federal Rule 506(c)', 'Ilimitados acreditados, verificación', '$3K-$8K', 'Fase 3'],
        ],

        'projection_title': 'Proyección Financiera — 12 Meses',
        'projection_subtitle': 'Escenario conservador: 70% reinversión, 8% default, gastos operativos incluidos.',
        'projection_table': [
            ['Mes', 'Capital ($27K)', 'Ganancia Acum.', 'Capital ($80K)', 'Ganancia Acum.', 'Clientes'],
            ['Mes 1', '$27,000', '$2,500', '$80,000', '$7,000', '40+'],
            ['Mes 3', '$33,000', '$8,500', '$100,000', '$25,000', '80+'],
            ['Mes 6', '$45,000', '$22,000', '$140,000', '$65,000', '150+'],
            ['Mes 9', '$60,000', '$40,000', '$190,000', '$120,000', '250+'],
            ['Mes 12', '$80,000', '$65,000', '$250,000+', '$190,000+', '400+'],
        ],
        'risk_title': '⚠️ Gestión de Riesgos',
        'risk_table': [
            ['Riesgo', 'Probabilidad', 'Impacto', 'Mitigación'],
            ['Default de clientes', '15%', 'Alto', 'Diversificar en 50+ préstamos, verificación de ingresos'],
            ['Cambio regulatorio', '5%', 'Alto', 'Operar dentro de límites OCCC, asesoría legal'],
            ['Falta de capital', '20%', 'Medio', 'Modelo de inversionistas, reinversión, reserva'],
        ],

        'action_title': 'Plan de Acción — 90 Días',
        'phases': [
            {
                'label': 'FASE 1',
                'title': 'Fundación\n(Semanas 1-4)',
                'items': [
                    'Finalizar licencia OCCC',
                    'Contratos legales listos',
                    'Cuenta bancaria empresarial',
                    'Procesador de pagos ACH',
                    'CRM y portal de clientes',
                    'Primer préstamo de prueba',
                ],
            },
            {
                'label': 'FASE 2',
                'title': 'Operación\n(Semanas 5-8)',
                'items': [
                    'Lanzar 20-30 préstamos Sub F',
                    'Activar cobros automatizados',
                    'Marketing comunitario local',
                    'Campaña de referidos',
                    'Evaluar tasa de default real',
                    'Reportar a credit bureau',
                ],
            },
            {
                'label': 'FASE 3',
                'title': 'Escalamiento\n(Semanas 9-12)',
                'items': [
                    'Introducir préstamos Sub E',
                    'Captar primer inversionista',
                    'Evaluar resultados, ajustar',
                    'Documentar historial',
                    'Preparar ronda de inversión',
                    'Expandir a ciudades cercanas',
                ],
            },
        ],

        'closing_title': '¿Listo para Invertir?',
        'closing_license': 'Texas OCCC Regulated Lender License | Chapter 342 Texas Finance Code',
    },

    'en': {
        'cover_subtitle': 'Comprehensive Business Plan 2026-2028',
        'cover_tagline': 'Smart Loans. Real Results. | Your Financial Partner Every Step of the Way.',
        'confidential': 'CONFIDENTIAL — For internal use only',
        'footer': 'Ross Lending Solutions LLC | 305 Bruce Ave, Dumas, TX 79029 | CONFIDENTIAL',

        'exec_title': 'Executive Summary',
        'exec_cards': [
            {
                'icon': '🏢',
                'title': 'The Company',
                'bullets': [
                    'Texas-regulated lender under OCCC license Ch. 342',
                    'Short and medium-term personal loans',
                    'Focus: Hispanic community in Texas Panhandle',
                    'Fully bilingual service (EN/ES)',
                    '24-hour approval process',
                ]
            },
            {
                'icon': '🎯',
                'title': 'Market Opportunity',
                'bullets': [
                    '65% of Panhandle population is Hispanic',
                    'Limited access to traditional bank credit',
                    '$2-5M annual demand in micro-loans',
                    'Low regulated competition in the area',
                    'Existing base of 900+ clients (Ross Tax)',
                ]
            },
            {
                'icon': '💰',
                'title': 'Investment Proposition',
                'bullets': [
                    'Starting capital: $27,000 - $80,000',
                    'Projected ROI: 180-240% annually',
                    'Investor return: 8-15% annually',
                    'Diversified loans (50+ micro-loans)',
                    'Regulated and state-compliant',
                ]
            },
        ],

        'reg_title': 'Regulatory Framework',
        'reg_subtitle': 'OCCC License — Texas Finance Code, Chapter 342 (One license covers both subchapters)',
        'reg_table': [
            ['Subchapter', 'Legal Section', 'Type', 'Amounts', 'Terms', 'Max Rate'],
            ['Subchapter F', '§342.251-259', 'Short-term', '$100 - $1,500', '1-3 months', '180-240% APR'],
            ['Subchapter E', '§342.201-214', 'Installment', '$2,000 - $12,000', '6-48 months', '18-30% APR'],
        ],
        'compliance_title': 'Compliance',
        'compliance_items': [
            'Truth in Lending Act (TILA) — Mandatory APR and fee disclosure',
            'Regulation Z — Standard disclosure format',
            'Equal Credit Opportunity Act — Non-discrimination',
            'FDCPA & Texas Debt Collection Act — Fair collection practices',
        ],

        'products_title': 'Loan Products',
        'subf_title': '⚡ Quick Loan — Subchapter F',
        'subf_table': [
            ['Amount', 'APR', 'Term', 'Charge', 'Total'],
            ['$200', '240%', '1 month', '$40', '$240'],
            ['$500', '180%', '1 month', '$75', '$575'],
            ['$1,000', '180%', '2 months', '$170', '$1,170'],
            ['$1,500', '180%', '3 months', '$240', '$1,740'],
        ],
        'sube_title': '📊 Installment Loan — Subchapter E',
        'sube_table': [
            ['Amount', 'Rate', 'Term', 'Monthly', 'Total'],
            ['$2,000', '~28%', '12 months', '$130', '$2,425'],
            ['$5,000', '~21%', '24 months', '$255', '$6,245'],
            ['$8,000', '~20%', '36 months', '$297', '$10,825'],
            ['$12,000', '~18%', '48 months', '$353', '$16,940'],
        ],
        'tax_advance_title': '🏛️ Tax Advance',
        'tax_advance_desc': 'Bridge loan (Jan-Apr) automatically repaid from IRS refund. Amounts: $200-$1,800 | Sub F rates | Base of 900+ tax season clients.',

        'capital_title': 'Capital Strategy',
        'scenario_27k': '📋 Scenario: $27,000 (Hybrid)',
        'table_27k': [
            ['Component', 'Capital', 'Loans/Mo', 'Profit'],
            ['Sub F — 70%', '$18,900', '~38', '$2,850/mo'],
            ['Sub E — 30%', '$8,100', '~3', '$160/mo'],
            ['TOTAL', '$27,000', '~41', '$3,010+/mo'],
        ],
        'scenario_80k': '🚀 Scenario: $80,000 (Scale-up)',
        'table_80k': [
            ['Component', 'Capital', 'Loans/Mo', 'Profit'],
            ['Sub F Mix', '$48,000', '~96', '$7,688/mo'],
            ['Sub E Mix', '$24,000', '~6', '$362/mo'],
            ['TOTAL', '$80,000', '~100+', '$8,050+/mo'],
        ],
        'key_metrics': [
            {'value': '$3,010+', 'label': 'Monthly Income\n(with $27K)'},
            {'value': '$8,050+', 'label': 'Monthly Income\n(with $80K)'},
            {'value': '180-240%', 'label': 'Annual ROI\nProjected'},
            {'value': '8-15%', 'label': 'Annual Return\nto Investor'},
        ],

        'investor_title': 'Investor Model',
        'investor_table': [
            ['Note Type', 'Annual Return', 'Risk', 'Min Investment', 'Interest Payments'],
            ['Private Note', '8-10%', 'Low', '$10,000', 'Monthly or Quarterly'],
            ['Participating Note', '10-12%', 'Moderate', '$25,000', 'Quarterly'],
            ['Profit Sharing', '12-15%', 'Medium', '$50,000', 'Monthly'],
            ['Partnership', '15-20%', 'High', '$100,000', 'Monthly'],
        ],
        'sec_title': 'SEC / Texas Securities Exemptions',
        'sec_table': [
            ['Exemption', 'Requirements', 'Legal Cost', 'Phase'],
            ['Texas §4005.012(a)(2)', 'Max 15 investors, sophisticated, TX only', '$0', 'Phase 1'],
            ['Federal Rule 506(b)', 'Max 35 non-accredited, PPM required', '$2K-$5K', 'Phase 2'],
            ['Federal Rule 506(c)', 'Unlimited accredited, verification', '$3K-$8K', 'Phase 3'],
        ],

        'projection_title': 'Financial Projections — 12 Months',
        'projection_subtitle': 'Conservative scenario: 70% reinvestment, 8% default, operating expenses included.',
        'projection_table': [
            ['Month', 'Capital ($27K)', 'Accum. Profit', 'Capital ($80K)', 'Accum. Profit', 'Clients'],
            ['Month 1', '$27,000', '$2,500', '$80,000', '$7,000', '40+'],
            ['Month 3', '$33,000', '$8,500', '$100,000', '$25,000', '80+'],
            ['Month 6', '$45,000', '$22,000', '$140,000', '$65,000', '150+'],
            ['Month 9', '$60,000', '$40,000', '$190,000', '$120,000', '250+'],
            ['Month 12', '$80,000', '$65,000', '$250,000+', '$190,000+', '400+'],
        ],
        'risk_title': '⚠️ Risk Management',
        'risk_table': [
            ['Risk', 'Probability', 'Impact', 'Mitigation'],
            ['Client default', '15%', 'High', 'Diversify into 50+ loans, income verification'],
            ['Regulatory change', '5%', 'High', 'Operate within OCCC limits, legal counsel'],
            ['Capital shortage', '20%', 'Medium', 'Investor model, reinvestment, reserve fund'],
        ],

        'action_title': '90-Day Action Plan',
        'phases': [
            {
                'label': 'PHASE 1',
                'title': 'Foundation\n(Weeks 1-4)',
                'items': [
                    'Finalize OCCC license',
                    'Legal contracts ready',
                    'Business bank account',
                    'ACH payment processor',
                    'CRM and client portal',
                    'First test loan',
                ],
            },
            {
                'label': 'PHASE 2',
                'title': 'Operations\n(Weeks 5-8)',
                'items': [
                    'Launch 20-30 Sub F loans',
                    'Activate auto-collections',
                    'Local community marketing',
                    'Referral campaign',
                    'Evaluate real default rate',
                    'Report to credit bureau',
                ],
            },
            {
                'label': 'PHASE 3',
                'title': 'Scaling\n(Weeks 9-12)',
                'items': [
                    'Introduce Sub E loans',
                    'Secure first investor',
                    'Evaluate results, adjust',
                    'Document track record',
                    'Prepare investment round',
                    'Expand to nearby cities',
                ],
            },
        ],

        'closing_title': 'Ready to Invest?',
        'closing_license': 'Texas OCCC Regulated Lender License | Chapter 342 Texas Finance Code',
    },
}


if __name__ == "__main__":
    os.makedirs('/app/backend/generated_docs', exist_ok=True)

    for lang in ['es', 'en']:
        prs = build_presentation(lang)
        filename = f"RLS_Pitch_Deck_{'ES' if lang == 'es' else 'EN'}_{datetime.now().strftime('%Y%m%d')}.pptx"
        filepath = f"/app/backend/generated_docs/{filename}"
        prs.save(filepath)
        size_mb = os.path.getsize(filepath) / (1024 * 1024)
        print(f"✅ {filename} — {size_mb:.1f} MB")

    print("\n📊 Pitch decks generated successfully!")
