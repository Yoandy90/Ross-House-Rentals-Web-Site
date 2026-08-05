"""
Script para generar una presentación PowerPoint profesional
sobre universidades de Astronomía y las mejores/más caras de EEUU.
Enviar por email a yoandyross@gmail.com
"""
import os
import io
import base64
import httpx
import smtplib
from email import encoders
from email.mime.base import MIMEBase
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from dotenv import load_dotenv

load_dotenv()

# ═══════════════════════════════════
# COLOR PALETTE
# ═══════════════════════════════════
DARK_NAVY = RGBColor(0x0D, 0x1B, 0x2A)
NAVY = RGBColor(0x1B, 0x26, 0x3B)
BLUE_GRAY = RGBColor(0x41, 0x5A, 0x77)
LIGHT_BLUE = RGBColor(0x77, 0x8D, 0xA9)
ACCENT_GOLD = RGBColor(0xFF, 0xB7, 0x03)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x88, 0x88, 0x88)
SOFT_BG = RGBColor(0xF8, 0xF9, 0xFA)
RED_ACCENT = RGBColor(0xE7, 0x4C, 0x3C)
TEAL = RGBColor(0x00, 0x96, 0x88)

# Slide dimensions: Widescreen 16:9
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=16, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(4)):
    """Add a new paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


# ═══════════════════════════════════
# UNIVERSITY DATA
# ═══════════════════════════════════

astronomy_unis = [
    {"name": "MIT", "loc": "Cambridge, MA", "tuition": "$61,990", "accept": "3.9%", "rank": "#1 QS 2025", "type": "Privada",
     "highlights": "Kavli Institute, NASA labs, 93% recibe ayuda financiera", "total4": "~$248K"},
    {"name": "Harvard University", "loc": "Cambridge, MA", "tuition": "$59,076", "accept": "3.6%", "rank": "#2 QS 2025", "type": "Privada",
     "highlights": "Harvard-Smithsonian CfA, observatorio propio", "total4": "~$236K"},
    {"name": "Caltech", "loc": "Pasadena, CA", "tuition": "$63,402", "accept": "3.1%", "rank": "#3 EduRank", "type": "Privada",
     "highlights": "Observatorios Palomar & Keck, NASA JPL, ratio 3:1", "total4": "~$254K"},
    {"name": "Princeton University", "loc": "Princeton, NJ", "tuition": "$59,710", "accept": "5.7%", "rank": "Top 5", "type": "Privada",
     "highlights": "Depto. exclusivo de Astrofísica, proyecto WFIRST NASA", "total4": "~$239K"},
    {"name": "Stanford University", "loc": "Stanford, CA", "tuition": "$62,484", "accept": "3.7%", "rank": "Top 5 QS", "type": "Privada",
     "highlights": "KIPAC, SLAC Lab, Silicon Valley tech espacial", "total4": "~$250K"},
    {"name": "UC Berkeley", "loc": "Berkeley, CA", "tuition": "$14,850*", "accept": "11.6%", "rank": "#1 Pública", "type": "Pública",
     "highlights": "MEJOR VALOR, Observatorio Lick, programa SETI", "total4": "~$59K*"},
    {"name": "UCLA", "loc": "Los Angeles, CA", "tuition": "$11,678*", "accept": "8.6%", "rank": "Top 10", "type": "Pública",
     "highlights": "Nobel Andrea Ghez, agujeros negros, matrícula baja", "total4": "~$47K*"},
    {"name": "U. of Chicago", "loc": "Chicago, IL", "tuition": "$66,939", "accept": "5.2%", "rank": "Top 5", "type": "Privada",
     "highlights": "Observatorio Yerkes, Fermi Lab, cosmología teórica", "total4": "~$268K"},
    {"name": "Columbia University", "loc": "New York, NY", "tuition": "$69,045", "accept": "3.9%", "rank": "Top 10", "type": "Privada",
     "highlights": "Hayden Planetarium, LIGO, Am. Museum Nat. History", "total4": "~$276K"},
    {"name": "Cornell University", "loc": "Ithaca, NY", "tuition": "$65,204", "accept": "7.9%", "rank": "Top 15", "type": "Privada",
     "highlights": "Legado Carl Sagan, James Webb Telescope, planetaria", "total4": "~$261K"},
    {"name": "U. of Arizona", "loc": "Tucson, AZ", "tuition": "$12,950*", "accept": "87%", "rank": "Top 5 Inv.", "type": "Pública",
     "highlights": "MAS ACCESIBLE, Steward Obs., misión OSIRIS-REx", "total4": "~$52K*"},
    {"name": "CU Boulder", "loc": "Boulder, CO", "tuition": "$13,194*", "accept": "80%", "rank": "Top 10 Planet.", "type": "Pública",
     "highlights": "LASP, NASA/NOAA, ingeniería aeroespacial", "total4": "~$53K*"},
    {"name": "UH Mānoa", "loc": "Honolulu, HI", "tuition": "$11,520*", "accept": "70%", "rank": "Top Observación", "type": "Pública",
     "highlights": "Mauna Kea (mejor sitio del mundo), 13 telescopios", "total4": "~$46K*"},
    {"name": "UC Santa Cruz", "loc": "Santa Cruz, CA", "tuition": "$14,100*", "accept": "47%", "rank": "Top 15", "type": "Pública",
     "highlights": "UCO/Lick Observatory, vistas al Pacífico", "total4": "~$56K*"},
    {"name": "UIUC", "loc": "Champaign, IL", "tuition": "$18,372*", "accept": "45%", "rank": "Top 20", "type": "Pública",
     "highlights": "NCSA supercomputación, astronomía computacional", "total4": "~$74K*"},
]

best_expensive = [
    {"name": "Columbia University", "tuition": "$69,045", "total4": "~$276,000", "known_for": "Ivy League, NYC, artes y ciencias"},
    {"name": "U. of Chicago", "tuition": "$66,939", "total4": "~$268,000", "known_for": "Economía, física, Nobel prizes"},
    {"name": "USC", "tuition": "~$67,000", "total4": "~$268,000", "known_for": "Cine, negocios, ingeniería"},
    {"name": "Cornell University", "tuition": "$65,204", "total4": "~$261,000", "known_for": "Ivy League, veterinaria, ingeniería"},
    {"name": "Yale University", "tuition": "~$64,700", "total4": "~$259,000", "known_for": "Derecho, artes, humanidades"},
    {"name": "Caltech", "tuition": "$63,402", "total4": "~$254,000", "known_for": "Ciencias, ingeniería, NASA JPL"},
    {"name": "Stanford", "tuition": "$62,484", "total4": "~$250,000", "known_for": "Tech, negocios, medicina"},
    {"name": "MIT", "tuition": "$61,990", "total4": "~$248,000", "known_for": "Ingeniería, tech, ciencias"},
]

best_overall = [
    {"rank": "1", "name": "MIT", "loc": "Cambridge, MA", "strong": "Ingeniería, Ciencias, Tech"},
    {"rank": "2", "name": "Stanford", "loc": "Stanford, CA", "strong": "Tech, Negocios, Medicina"},
    {"rank": "3", "name": "Harvard", "loc": "Cambridge, MA", "strong": "Derecho, Medicina, Ciencias"},
    {"rank": "4", "name": "Caltech", "loc": "Pasadena, CA", "strong": "Ciencias, Ingeniería, Astro"},
    {"rank": "5", "name": "Princeton", "loc": "Princeton, NJ", "strong": "Matemáticas, Física, Econ"},
    {"rank": "6", "name": "U. of Chicago", "loc": "Chicago, IL", "strong": "Economía, Física, Derecho"},
    {"rank": "7", "name": "Yale", "loc": "New Haven, CT", "strong": "Derecho, Artes, Humanidades"},
    {"rank": "8", "name": "Columbia", "loc": "New York, NY", "strong": "Periodismo, Negocios, Artes"},
    {"rank": "9", "name": "UPenn", "loc": "Philadelphia, PA", "strong": "Negocios (Wharton), Medicina"},
    {"rank": "10", "name": "Duke", "loc": "Durham, NC", "strong": "Medicina, Derecho, Ciencias"},
]


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ════════════════════════════════════════
    # SLIDE 1: TITLE
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, DARK_NAVY)

    # Top accent line
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT_GOLD)

    # Title
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                 "UNIVERSIDADES EN ESTADOS UNIDOS", 44, WHITE, True, PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
                 "Las Mejores  |  Las Más Caras  |  Astronomía & Astrofísica", 26, ACCENT_GOLD, False, PP_ALIGN.CENTER)

    # Divider
    add_shape(slide, Inches(4.5), Inches(3.8), Inches(4), Inches(0.04), ACCENT_GOLD)

    # Bottom info
    add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
                 "Guía Completa para la Familia Ross", 22, LIGHT_BLUE, False, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
                 "Febrero 2026", 18, MED_GRAY, False, PP_ALIGN.CENTER)

    # Bottom accent
    add_shape(slide, Inches(0), Inches(7.42), SLIDE_WIDTH, Inches(0.08), ACCENT_GOLD)

    # ════════════════════════════════════════
    # SLIDE 2: AGENDA
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), DARK_NAVY)
    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
                 "CONTENIDO", 32, WHITE, True, PP_ALIGN.CENTER)

    sections = [
        ("01", "Las 10 Mejores Universidades de EEUU", "Rankings generales 2025-2026"),
        ("02", "Las Universidades Más Caras", "Comparativa de costos anuales y totales"),
        ("03", "Top 15 para Astronomía & Astrofísica", "Programas especializados en astronomía"),
        ("04", "Tabla Comparativa de Costos (Astronomía)", "Matrícula residente vs no-residente"),
        ("05", "Mejores Opciones por Categoría", "Recomendaciones según prioridades"),
        ("06", "Próximos Pasos", "Plan de acción sugerido"),
    ]

    y = 1.5
    for num, title, desc in sections:
        # Number circle
        shape = add_shape(slide, Inches(1.5), Inches(y), Inches(0.7), Inches(0.7), DARK_NAVY)
        shape.text_frame.paragraphs[0].text = num
        shape.text_frame.paragraphs[0].font.size = Pt(20)
        shape.text_frame.paragraphs[0].font.color.rgb = ACCENT_GOLD
        shape.text_frame.paragraphs[0].font.bold = True
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        shape.text_frame.word_wrap = False

        add_text_box(slide, Inches(2.5), Inches(y - 0.05), Inches(8), Inches(0.4),
                     title, 20, DARK_NAVY, True)
        add_text_box(slide, Inches(2.5), Inches(y + 0.35), Inches(8), Inches(0.35),
                     desc, 14, MED_GRAY, False)
        y += 0.9

    # ════════════════════════════════════════
    # SLIDE 3: TOP 10 BEST UNIVERSITIES
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), DARK_NAVY)
    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
                 "LAS 10 MEJORES UNIVERSIDADES DE EEUU", 30, WHITE, True, PP_ALIGN.CENTER)

    # Table header
    header_y = 1.4
    add_shape(slide, Inches(0.8), Inches(header_y), Inches(11.7), Inches(0.5), NAVY)
    headers = [("#", 0.8, 0.6), ("Universidad", 1.5, 3.5), ("Ubicación", 5.2, 2.8), ("Fortalezas", 8.2, 4.3)]
    for text, x, w in headers:
        add_text_box(slide, Inches(x), Inches(header_y + 0.05), Inches(w), Inches(0.4),
                     text, 14, ACCENT_GOLD, True, PP_ALIGN.CENTER)

    # Rows
    row_y = header_y + 0.55
    for i, u in enumerate(best_overall):
        bg = SOFT_BG if i % 2 == 0 else WHITE
        add_shape(slide, Inches(0.8), Inches(row_y), Inches(11.7), Inches(0.5), bg)

        # Rank with special color for top 3
        rank_color = ACCENT_GOLD if int(u["rank"]) <= 3 else DARK_TEXT
        add_text_box(slide, Inches(0.8), Inches(row_y + 0.05), Inches(0.6), Inches(0.4),
                     u["rank"], 16, rank_color, True, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(1.5), Inches(row_y + 0.05), Inches(3.5), Inches(0.4),
                     u["name"], 15, DARK_NAVY, True)
        add_text_box(slide, Inches(5.2), Inches(row_y + 0.05), Inches(2.8), Inches(0.4),
                     u["loc"], 14, DARK_TEXT, False)
        add_text_box(slide, Inches(8.2), Inches(row_y + 0.05), Inches(4.3), Inches(0.4),
                     u["strong"], 13, BLUE_GRAY, False)
        row_y += 0.52

    add_text_box(slide, Inches(0.8), Inches(row_y + 0.2), Inches(11), Inches(0.4),
                 "Fuente: QS World University Rankings, US News & World Report 2025-2026", 11, MED_GRAY, False, PP_ALIGN.CENTER)

    # ════════════════════════════════════════
    # SLIDE 4: MOST EXPENSIVE
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), RGBColor(0x8B, 0x00, 0x00))
    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
                 "LAS UNIVERSIDADES MÁS CARAS DE EEUU", 30, WHITE, True, PP_ALIGN.CENTER)

    # Table
    header_y = 1.4
    add_shape(slide, Inches(0.8), Inches(header_y), Inches(11.7), Inches(0.5), RGBColor(0x8B, 0x00, 0x00))
    exp_headers = [("Universidad", 0.8, 2.8), ("Matrícula/Año", 3.8, 2.0), ("Costo Total 4 Años", 6.0, 2.2), ("Conocida Por", 8.4, 4.1)]
    for text, x, w in exp_headers:
        add_text_box(slide, Inches(x), Inches(header_y + 0.05), Inches(w), Inches(0.4),
                     text, 14, WHITE, True, PP_ALIGN.CENTER)

    row_y = header_y + 0.55
    for i, u in enumerate(best_expensive):
        bg = RGBColor(0xFF, 0xF0, 0xF0) if i % 2 == 0 else WHITE
        add_shape(slide, Inches(0.8), Inches(row_y), Inches(11.7), Inches(0.5), bg)
        add_text_box(slide, Inches(0.8), Inches(row_y + 0.05), Inches(2.8), Inches(0.4),
                     u["name"], 15, DARK_NAVY, True)
        add_text_box(slide, Inches(3.8), Inches(row_y + 0.05), Inches(2.0), Inches(0.4),
                     u["tuition"], 15, RED_ACCENT, True, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(6.0), Inches(row_y + 0.05), Inches(2.2), Inches(0.4),
                     u["total4"], 14, DARK_TEXT, False, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(8.4), Inches(row_y + 0.05), Inches(4.1), Inches(0.4),
                     u["known_for"], 13, BLUE_GRAY, False)
        row_y += 0.52

    # Important note
    note_shape = add_shape(slide, Inches(1.5), Inches(row_y + 0.3), Inches(10), Inches(0.8), RGBColor(0xFF, 0xF8, 0xE1), ACCENT_GOLD)
    note_shape.text_frame.word_wrap = True
    note_shape.text_frame.paragraphs[0].text = "IMPORTANTE: La mayoría de estas universidades ofrecen ayuda financiera generosa."
    note_shape.text_frame.paragraphs[0].font.size = Pt(13)
    note_shape.text_frame.paragraphs[0].font.bold = True
    note_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x8B, 0x6A, 0x00)
    p2 = note_shape.text_frame.add_paragraph()
    p2.text = "Harvard/MIT/Princeton: Familias con ingresos < $100K/año pueden estudiar GRATIS."
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(0x8B, 0x6A, 0x00)

    # ════════════════════════════════════════
    # SLIDE 5: ASTRONOMY TOP 15 - PRIVATE
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), DARK_NAVY)
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.4),
                 "TOP UNIVERSIDADES PARA ASTRONOMÍA", 28, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(0.55), Inches(12), Inches(0.4),
                 "Universidades Privadas", 18, ACCENT_GOLD, False, PP_ALIGN.CENTER)

    # Table header
    header_y = 1.3
    add_shape(slide, Inches(0.3), Inches(header_y), Inches(12.7), Inches(0.45), NAVY)
    astro_h = [("Universidad", 0.3, 2.2), ("Ubicación", 2.6, 1.8), ("Matrícula/Año", 4.5, 1.5), ("Aceptación", 6.1, 1.2),
               ("Costo 4 Años", 7.4, 1.3), ("Ranking", 8.8, 1.2), ("Puntos Destacados", 10.1, 2.9)]
    for text, x, w in astro_h:
        add_text_box(slide, Inches(x), Inches(header_y + 0.03), Inches(w), Inches(0.4),
                     text, 11, ACCENT_GOLD, True, PP_ALIGN.CENTER)

    row_y = header_y + 0.48
    private_unis = [u for u in astronomy_unis if u["type"] == "Privada"]
    for i, u in enumerate(private_unis):
        bg = SOFT_BG if i % 2 == 0 else WHITE
        add_shape(slide, Inches(0.3), Inches(row_y), Inches(12.7), Inches(0.55), bg)
        add_text_box(slide, Inches(0.3), Inches(row_y + 0.05), Inches(2.2), Inches(0.4),
                     u["name"], 12, DARK_NAVY, True)
        add_text_box(slide, Inches(2.6), Inches(row_y + 0.05), Inches(1.8), Inches(0.4),
                     u["loc"], 11, DARK_TEXT)
        add_text_box(slide, Inches(4.5), Inches(row_y + 0.05), Inches(1.5), Inches(0.4),
                     u["tuition"], 12, RED_ACCENT, True, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(6.1), Inches(row_y + 0.05), Inches(1.2), Inches(0.4),
                     u["accept"], 12, DARK_TEXT, False, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(7.4), Inches(row_y + 0.05), Inches(1.3), Inches(0.4),
                     u["total4"], 11, DARK_TEXT, False, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(8.8), Inches(row_y + 0.05), Inches(1.2), Inches(0.4),
                     u["rank"], 10, TEAL, True, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(10.1), Inches(row_y + 0.05), Inches(2.9), Inches(0.45),
                     u["highlights"], 9, BLUE_GRAY)
        row_y += 0.58

    # ════════════════════════════════════════
    # SLIDE 6: ASTRONOMY - PUBLIC
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), TEAL)
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.4),
                 "TOP UNIVERSIDADES PARA ASTRONOMÍA", 28, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(0.55), Inches(12), Inches(0.4),
                 "Universidades Públicas — Mejor Valor", 18, RGBColor(0xB2, 0xFF, 0xD8), False, PP_ALIGN.CENTER)

    header_y = 1.3
    add_shape(slide, Inches(0.3), Inches(header_y), Inches(12.7), Inches(0.45), RGBColor(0x00, 0x6B, 0x5C))
    for text, x, w in astro_h:
        add_text_box(slide, Inches(x), Inches(header_y + 0.03), Inches(w), Inches(0.4),
                     text, 11, WHITE, True, PP_ALIGN.CENTER)

    row_y = header_y + 0.48
    public_unis = [u for u in astronomy_unis if u["type"] == "Pública"]
    for i, u in enumerate(public_unis):
        bg = RGBColor(0xE8, 0xF5, 0xE9) if i % 2 == 0 else WHITE
        add_shape(slide, Inches(0.3), Inches(row_y), Inches(12.7), Inches(0.55), bg)
        add_text_box(slide, Inches(0.3), Inches(row_y + 0.05), Inches(2.2), Inches(0.4),
                     u["name"], 12, DARK_NAVY, True)
        add_text_box(slide, Inches(2.6), Inches(row_y + 0.05), Inches(1.8), Inches(0.4),
                     u["loc"], 11, DARK_TEXT)
        add_text_box(slide, Inches(4.5), Inches(row_y + 0.05), Inches(1.5), Inches(0.4),
                     u["tuition"], 12, ACCENT_GREEN, True, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(6.1), Inches(row_y + 0.05), Inches(1.2), Inches(0.4),
                     u["accept"], 12, DARK_TEXT, False, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(7.4), Inches(row_y + 0.05), Inches(1.3), Inches(0.4),
                     u["total4"], 11, DARK_TEXT, False, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(8.8), Inches(row_y + 0.05), Inches(1.2), Inches(0.4),
                     u["rank"], 10, TEAL, True, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(10.1), Inches(row_y + 0.05), Inches(2.9), Inches(0.45),
                     u["highlights"], 9, BLUE_GRAY)
        row_y += 0.58

    # Note at bottom
    add_text_box(slide, Inches(0.8), Inches(row_y + 0.15), Inches(11), Inches(0.4),
                 "* Matrícula para residentes del estado. No-residentes pagan ~$33K-$45K/año.", 12, MED_GRAY, False, PP_ALIGN.CENTER)

    # ════════════════════════════════════════
    # SLIDE 7: COST COMPARISON VISUAL
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), DARK_NAVY)
    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
                 "COMPARATIVA DE COSTOS: 4 AÑOS DE CARRERA", 28, WHITE, True, PP_ALIGN.CENTER)

    # Visual bars
    bar_data = [
        ("Columbia", 276, RED_ACCENT, "Privada"),
        ("U. Chicago", 268, RED_ACCENT, "Privada"),
        ("Cornell", 261, RED_ACCENT, "Privada"),
        ("Caltech", 254, RED_ACCENT, "Privada"),
        ("Stanford", 250, RGBColor(0xE0, 0x6C, 0x5C), "Privada"),
        ("MIT", 248, RGBColor(0xE0, 0x6C, 0x5C), "Privada"),
        ("Princeton", 239, RGBColor(0xE0, 0x8C, 0x5C), "Privada"),
        ("Harvard", 236, RGBColor(0xE0, 0x8C, 0x5C), "Privada"),
        ("UIUC*", 74, ACCENT_GREEN, "Pública"),
        ("UC Berkeley*", 59, ACCENT_GREEN, "Pública"),
        ("UC Santa Cruz*", 56, ACCENT_GREEN, "Pública"),
        ("CU Boulder*", 53, ACCENT_GREEN, "Pública"),
        ("U. Arizona*", 52, ACCENT_GREEN, "Pública"),
        ("UCLA*", 47, TEAL, "Pública"),
        ("UH Mānoa*", 46, TEAL, "Pública"),
    ]

    y = 1.4
    max_cost = 280
    bar_max_width = 7.5

    for name, cost, color, utype in bar_data:
        bar_width = (cost / max_cost) * bar_max_width
        add_text_box(slide, Inches(0.3), Inches(y - 0.02), Inches(2.2), Inches(0.35),
                     name, 11, DARK_NAVY, True, PP_ALIGN.RIGHT)
        add_shape(slide, Inches(2.6), Inches(y), Inches(bar_width), Inches(0.3), color)
        add_text_box(slide, Inches(2.6 + bar_width + 0.1), Inches(y - 0.02), Inches(1.5), Inches(0.35),
                     f"~${cost}K", 11, DARK_TEXT, True)
        y += 0.37

    add_text_box(slide, Inches(0.8), Inches(y + 0.15), Inches(11), Inches(0.4),
                 "* Costo para residentes del estado  |  Rojo = Privada  |  Verde = Pública", 11, MED_GRAY, False, PP_ALIGN.CENTER)

    # ════════════════════════════════════════
    # SLIDE 8: RECOMMENDATIONS BY CATEGORY
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), DARK_NAVY)
    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
                 "RECOMENDACIONES POR CATEGORÍA", 30, WHITE, True, PP_ALIGN.CENTER)

    recommendations = [
        ("Mejor Programa General", "MIT / Caltech", "Lideran investigación astronómica global", DARK_NAVY),
        ("Mejor Valor (Pública)", "UC Berkeley", "Calidad élite con matrícula pública ~$15K/año", TEAL),
        ("Más Accesible + Top Inv.", "U. of Arizona", "87% aceptación, Steward Observatory, NASA", ACCENT_GREEN),
        ("Matrícula Más Baja", "UCLA / UH Mānoa", "$11,520-$11,678/año para residentes", RGBColor(0x27, 0xAE, 0x60)),
        ("Mejor para Observación", "UH Mānoa", "Mauna Kea: el mejor sitio de observación del mundo", BLUE_GRAY),
        ("Mejor Ciencia Planetaria", "U. Arizona / Cornell", "Misiones NASA, legado Carl Sagan", RGBColor(0x8E, 0x44, 0xAD)),
    ]

    # Cards layout: 3x2
    card_w = 3.5
    card_h = 1.6
    start_x = 0.7
    start_y = 1.5
    gap_x = 0.3
    gap_y = 0.3

    for idx, (cat, uni, reason, accent) in enumerate(recommendations):
        col = idx % 3
        row = idx // 3
        x = start_x + col * (card_w + gap_x)
        y_pos = start_y + row * (card_h + gap_y)

        # Card background
        card = add_shape(slide, Inches(x), Inches(y_pos), Inches(card_w), Inches(card_h), WHITE, RGBColor(0xDD, 0xDD, 0xDD))
        # Accent top bar
        add_shape(slide, Inches(x), Inches(y_pos), Inches(card_w), Inches(0.06), accent)

        add_text_box(slide, Inches(x + 0.15), Inches(y_pos + 0.15), Inches(card_w - 0.3), Inches(0.35),
                     cat, 13, accent, True)
        add_text_box(slide, Inches(x + 0.15), Inches(y_pos + 0.55), Inches(card_w - 0.3), Inches(0.4),
                     uni, 18, DARK_NAVY, True)
        add_text_box(slide, Inches(x + 0.15), Inches(y_pos + 1.0), Inches(card_w - 0.3), Inches(0.5),
                     reason, 11, MED_GRAY)

    # ════════════════════════════════════════
    # SLIDE 9: FINANCIAL AID INFO
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), RGBColor(0x1A, 0x5C, 0x2D))
    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
                 "DATO CLAVE: AYUDA FINANCIERA", 30, WHITE, True, PP_ALIGN.CENTER)

    # Big message
    msg_shape = add_shape(slide, Inches(1), Inches(1.5), Inches(11), Inches(2), RGBColor(0xE8, 0xF5, 0xE9))
    msg_shape.text_frame.word_wrap = True
    p = msg_shape.text_frame.paragraphs[0]
    p.text = "NO te dejes asustar por los precios"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x5C, 0x2D)
    p.alignment = PP_ALIGN.CENTER

    p2 = msg_shape.text_frame.add_paragraph()
    p2.text = "Las universidades privadas más caras son también las que más ayuda financiera dan"
    p2.font.size = Pt(18)
    p2.font.color.rgb = DARK_TEXT
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)

    # Facts
    facts = [
        ("Harvard", "Familias con ingresos < $85K/año = GRATIS. < $150K = muy reducida."),
        ("MIT", "El 93% de estudiantes recibe algún tipo de ayuda financiera."),
        ("Princeton", "100% basada en necesidad. No usan préstamos, solo becas."),
        ("Stanford", "Familias < $75K/año: matrícula + alojamiento = GRATIS."),
        ("Caltech", "Becas promedio de $45,000/año para estudiantes con necesidad."),
    ]

    y = 4.0
    for uni_name, fact in facts:
        add_shape(slide, Inches(1), Inches(y), Inches(2), Inches(0.45), DARK_NAVY)
        add_text_box(slide, Inches(1), Inches(y + 0.05), Inches(2), Inches(0.35),
                     uni_name, 14, WHITE, True, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(3.2), Inches(y + 0.05), Inches(8.5), Inches(0.35),
                     fact, 14, DARK_TEXT)
        y += 0.55

    # ════════════════════════════════════════
    # SLIDE 10: NEXT STEPS
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT_GOLD)

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7),
                 "PRÓXIMOS PASOS", 32, WHITE, True, PP_ALIGN.CENTER)
    add_shape(slide, Inches(4.5), Inches(1.3), Inches(4), Inches(0.04), ACCENT_GOLD)

    steps = [
        "Definir presupuesto familiar anual para educación universitaria",
        "Investigar programas de ayuda financiera y becas",
        "Visitar los campus de 3-5 universidades preferidas (tours virtuales disponibles)",
        "Preparar exámenes SAT/ACT (requeridos por la mayoría)",
        "Considerar programas de verano de astronomía para jóvenes",
        "Para universidades de California: investigar requisitos de residencia (1 año previo)",
        "Aplicar a múltiples universidades en diferentes niveles de selectividad",
    ]

    y = 1.8
    for i, step in enumerate(steps):
        # Number
        num_shape = add_shape(slide, Inches(1.5), Inches(y), Inches(0.55), Inches(0.55), ACCENT_GOLD)
        num_shape.text_frame.paragraphs[0].text = str(i + 1)
        num_shape.text_frame.paragraphs[0].font.size = Pt(18)
        num_shape.text_frame.paragraphs[0].font.color.rgb = DARK_NAVY
        num_shape.text_frame.paragraphs[0].font.bold = True
        num_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_text_box(slide, Inches(2.3), Inches(y + 0.08), Inches(9), Inches(0.45),
                     step, 16, WHITE)
        y += 0.7

    # Bottom
    add_text_box(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
                 "Preparado con datos actualizados a Febrero 2026", 12, MED_GRAY, False, PP_ALIGN.CENTER)
    add_shape(slide, Inches(0), Inches(7.42), SLIDE_WIDTH, Inches(0.08), ACCENT_GOLD)

    return prs


# ═══════════════════════════════════
# EMAIL SENDING
# ═══════════════════════════════════

def send_pptx_email(pptx_buffer, pdf_path, to_email="yoandyross@gmail.com"):
    """Send PPTX + PDF via SendGrid."""
    sendgrid_key = os.getenv("SENDGRID_API_KEY")
    from_email = os.getenv("SENDGRID_FROM_EMAIL", "info@rosstaxpreparation.com")

    if not sendgrid_key:
        print("ERROR: No SENDGRID_API_KEY")
        return False

    pptx_data = pptx_buffer.getvalue()
    pptx_b64 = base64.b64encode(pptx_data).decode("utf-8")

    attachments_list = [
        {
            "content": pptx_b64,
            "filename": "Universidades_EEUU_Presentacion_2026.pptx",
            "type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "disposition": "attachment",
        }
    ]

    # Also attach the PDF
    if os.path.exists(pdf_path):
        with open(pdf_path, "rb") as f:
            pdf_b64 = base64.b64encode(f.read()).decode("utf-8")
        attachments_list.append({
            "content": pdf_b64,
            "filename": "Universidades_Astronomia_EEUU_2026.pdf",
            "type": "application/pdf",
            "disposition": "attachment",
        })

    html_body = """
    <div style="font-family: 'Segoe UI', Arial, sans-serif; max-width: 600px; margin: 0 auto;">
        <div style="background: linear-gradient(135deg, #0D1B2A 0%, #1B263B 100%); padding: 30px; text-align: center; border-radius: 12px 12px 0 0;">
            <h1 style="color: white; margin: 0; font-size: 22px;">&#127891; Universidades en EEUU</h1>
            <p style="color: #FFB703; margin: 8px 0 0;">Presentación + PDF Completo</p>
        </div>
        <div style="background: #ffffff; padding: 30px; border: 1px solid #e0e0e0;">
            <p style="color: #333; font-size: 15px;">Hola Yoandy,</p>
            <p style="color: #333; font-size: 15px;">Adjunto encontrarás <strong>2 archivos</strong>:</p>
            <div style="background: #f0f4ff; padding: 15px; border-radius: 8px; margin: 15px 0;">
                <p style="margin: 5px 0; font-size: 14px;"><strong>📊 Presentación (.pptx)</strong> — 10 diapositivas para ver en TV</p>
                <ul style="font-size: 13px; color: #555;">
                    <li>Top 10 mejores universidades de EEUU</li>
                    <li>Las más caras con comparativa</li>
                    <li>Top 15 para Astronomía (privadas y públicas)</li>
                    <li>Gráfico visual de costos</li>
                    <li>Recomendaciones y próximos pasos</li>
                    <li>Info sobre ayuda financiera</li>
                </ul>
            </div>
            <div style="background: #fff3e0; padding: 15px; border-radius: 8px; margin: 15px 0;">
                <p style="margin: 5px 0; font-size: 14px;"><strong>📄 PDF Detallado</strong> — Perfiles completos de 15 universidades</p>
                <p style="font-size: 13px; color: #555; margin: 5px 0;">Direcciones, costos, rankings, highlights, y más.</p>
            </div>
            <div style="background: #e8f5e9; padding: 12px; border-radius: 8px; margin: 15px 0;">
                <p style="margin: 0; font-size: 13px; color: #1a5c2d;"><strong>💡 Para ver en el TV:</strong> Abre el archivo .pptx en PowerPoint o Google Slides y pon en modo presentación (F5).</p>
            </div>
        </div>
        <div style="background: #0D1B2A; padding: 15px; text-align: center; border-radius: 0 0 12px 12px;">
            <p style="color: rgba(255,255,255,0.7); font-size: 11px; margin: 0;">Ross Tax Preparation LLC — Febrero 2026</p>
        </div>
    </div>
    """

    payload = {
        "personalizations": [{"to": [{"email": to_email}]}],
        "from": {"email": from_email, "name": "Ross Tax Preparation"},
        "subject": "🎓 Presentación: Universidades de EEUU — Para ver en TV",
        "content": [
            {"type": "text/plain", "value": "Adjunto: Presentación de universidades para TV + PDF detallado."},
            {"type": "text/html", "value": html_body},
        ],
        "attachments": attachments_list,
    }

    try:
        resp = httpx.post(
            "https://api.sendgrid.com/v3/mail/send",
            json=payload,
            headers={
                "Authorization": f"Bearer {sendgrid_key}",
                "Content-Type": "application/json",
            },
            timeout=30.0,
        )
        if resp.status_code in (200, 201, 202):
            print(f"✅ Email enviado exitosamente a {to_email}")
            return True
        else:
            print(f"❌ Error SendGrid: {resp.status_code} - {resp.text}")
            return False
    except Exception as e:
        print(f"❌ Error: {e}")
        return False


# ═══════════════════════════════════
# MAIN
# ═══════════════════════════════════

if __name__ == "__main__":
    print("📊 Generando presentación PowerPoint...")
    prs = create_presentation()

    # Save to buffer
    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)

    # Save locally
    local_path = "/app/backend/static/Universidades_EEUU_Presentacion_2026.pptx"
    os.makedirs(os.path.dirname(local_path), exist_ok=True)
    with open(local_path, "wb") as f:
        f.write(pptx_buffer.getvalue())
    print(f"✅ Presentación guardada: {local_path}")

    pptx_buffer.seek(0)

    # PDF path from previous script
    pdf_path = "/app/backend/static/Universidades_Astronomia_EEUU_2026.pdf"

    print("📧 Enviando por email (PPTX + PDF)...")
    success = send_pptx_email(pptx_buffer, pdf_path)

    if success:
        print("\n🎉 ¡Todo listo! Presentación + PDF enviados a yoandyross@gmail.com")
    else:
        print("\n⚠️ Archivos generados localmente pero hubo error al enviar email.")
